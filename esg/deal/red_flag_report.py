"""
Red-flag report export.

Builds the deliverable and refuses to release it unless three things hold:

1. Every quantified exposure in it has been reviewed (esg.methodology.exposure).
2. Any benchmark it cites comes from a publishable cohort
   (esg.benchmarks.provenance) — illustrative peers cannot leave the building.
3. Required sign-off is complete against the exact bytes being exported
   (esg.deal.signoff).

The gates are applied at build time, so a report that would misrepresent its
basis is never produced, rather than produced and then policed.
"""

import io
import json

from esg import clock
from esg.benchmarks import provenance
from esg.db.scope import require_principal
from esg.deal import information_requests as irs
from esg.deal import signoff
from esg.methodology import exposure
from esg.security import audit, rbac

# PwC brand palette, from the deck template.
PWC_ORANGE = (0xFF, 0x5A, 0x00)
PWC_ROSE = (0xD0, 0x40, 0x72)
TEXT_DARK = (0x23, 0x23, 0x23)
TEXT_BODY = (0x46, 0x46, 0x46)


class ExportError(RuntimeError):
    pass


def build_payload(session, deal_id, company_id, findings, exposures=None,
                  benchmark_cohorts=None, reporting_year=None):
    """Assemble the report content and enforce the basis gates."""
    principal = require_principal()
    rbac.check(rbac.EXPORT_REPORT, deal_id=deal_id, principal=principal)

    exposures = exposures or []
    for run in exposures:
        # An unreviewed number cannot appear in a client document.
        exposure.require_review(run)

    for cohort in benchmark_cohorts or []:
        provenance.require_publishable(cohort, "the red-flag report")

    outstanding = irs.outstanding_at_signing(session, deal_id)

    quantified = [r for r in exposures
                  if exposure.present(r)["disclosure"] == exposure.QUANTIFIED]
    indicative = [r for r in exposures
                  if exposure.present(r)["disclosure"] == exposure.INDICATIVE]

    return {
        "deal_id": deal_id,
        "company_id": company_id,
        "reporting_year": reporting_year,
        "prepared_by": principal.username,
        "prepared_at": clock.now().isoformat(timespec="seconds"),
        "methodology_version": exposure.METHODOLOGY_VERSION,
        "findings": sorted(
            findings,
            key=lambda f: {"Critical": 0, "High": 1, "Medium": 2, "Low": 3}.get(
                f.get("severity"), 9
            ),
        ),
        "quantified_total_usd": round(
            sum(r.point_estimate_usd or 0 for r in quantified), 2
        ),
        "indicative_range_usd": [
            round(sum(r.low_usd or 0 for r in indicative), 2),
            round(sum(r.high_usd or 0 for r in indicative), 2),
        ] if indicative else None,
        "exposures": [
            {
                "finding_id": r.finding_id,
                "method": r.method,
                "presentation": exposure.present(r),
                "reviewed_by": r.reviewed_by,
                "methodology_version": r.methodology_version,
            }
            for r in exposures
        ],
        "benchmark_notes": [
            {"metric_code": c["metric_code"], "provenance": c["provenance"],
             "peer_count": c["peer_count"], "sources": c["sources"]}
            for c in (benchmark_cohorts or [])
        ],
        "outstanding_information_requests": outstanding,
    }


def render_docx(payload):
    """Word rendering of the report. Returns bytes."""
    import docx
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    document = docx.Document()

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = title.add_run("ESG Due Diligence — Red Flag Report")
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TEXT_DARK)

    subtitle = document.add_paragraph()
    run = subtitle.add_run(
        f"Deal {payload['deal_id']} · Target {payload['company_id']}"
        + (f" · FY{payload['reporting_year']}" if payload["reporting_year"] else "")
    )
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(*TEXT_BODY)

    meta = document.add_paragraph()
    run = meta.add_run(
        f"Prepared by {payload['prepared_by']} on {payload['prepared_at']} · "
        f"Exposure methodology {payload['methodology_version']}"
    )
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(*TEXT_BODY)

    document.add_heading("Exposure summary", level=1)
    document.add_paragraph(
        f"Quantified exposure (observed amounts): "
        f"USD {payload['quantified_total_usd']:,.0f}"
    )
    if payload["indicative_range_usd"]:
        low, high = payload["indicative_range_usd"]
        paragraph = document.add_paragraph(
            f"Indicative range (uncalibrated model): USD {low:,.0f} – {high:,.0f}"
        )
        run = paragraph.add_run(
            "\nThis range is not a quantified exposure. It derives from judgement "
            "parameters that have not been empirically calibrated, and is presented "
            "as a range because a point estimate would overstate the precision "
            "available."
        )
        run.font.size = Pt(9)
        run.italic = True

    document.add_heading("Findings", level=1)
    for index, finding in enumerate(payload["findings"], start=1):
        document.add_heading(
            f"{index}. [{finding.get('severity', 'Unrated')}] "
            f"{finding.get('finding') or finding.get('title', 'Finding')}",
            level=2,
        )
        if finding.get("why_it_matters"):
            document.add_paragraph(finding["why_it_matters"])
        citations = [
            c for c in (
                [finding.get("reported_citation"), finding.get("citation"),
                 finding.get("original_citation"), finding.get("restated_citation")]
                + list(finding.get("facility_citations") or [])
            ) if c
        ]
        if citations:
            paragraph = document.add_paragraph()
            run = paragraph.add_run("Evidence: " + "; ".join(dict.fromkeys(citations)))
            run.font.size = Pt(9)
            run.italic = True

    if payload["benchmark_notes"]:
        document.add_heading("Benchmark basis", level=1)
        for note in payload["benchmark_notes"]:
            document.add_paragraph(
                f"{note['metric_code']}: {note['peer_count']} peers, "
                f"provenance {note['provenance']}"
                + (f", sources: {', '.join(note['sources'][:3])}" if note["sources"] else "")
            )

    outstanding = payload["outstanding_information_requests"]
    document.add_heading("Outstanding information requests", level=1)
    if outstanding["count"]:
        for item in outstanding["items"]:
            document.add_paragraph(
                f"{item['reference']} · {item['title']} "
                f"({item['priority']}, {item['status']}, open {item['days_open']} days)",
                style="List Bullet",
            )
    document.add_paragraph(outstanding["note"])

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def render_pptx(payload):
    """Slide rendering for the deal-team readout. Returns bytes."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    def textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=TEXT_DARK):
        box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        # Formatting goes on the run, not the paragraph default, so it survives
        # rendering outside PowerPoint.
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        return box

    # Title slide
    slide = presentation.slides.add_slide(blank)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*PWC_ORANGE)
    bar.line.fill.background()
    textbox(slide, 0.7, 2.2, 11.5, 1.0,
            "ESG Due Diligence — Red Flag Report", size=40, bold=True)
    textbox(slide, 0.7, 3.3, 11.5, 0.6,
            f"Deal {payload['deal_id']} · Target {payload['company_id']}", size=18,
            color=TEXT_BODY)
    textbox(slide, 0.7, 6.4, 11.5, 0.4,
            f"Prepared {payload['prepared_at']} · Methodology "
            f"{payload['methodology_version']}", size=10, color=TEXT_BODY)

    # Exposure slide
    slide = presentation.slides.add_slide(blank)
    textbox(slide, 0.7, 0.4, 11.5, 0.6, "Exposure summary", size=28, bold=True)
    textbox(slide, 0.7, 1.4, 11.5, 0.5,
            f"Quantified (observed amounts): USD "
            f"{payload['quantified_total_usd']:,.0f}", size=18)
    if payload["indicative_range_usd"]:
        low, high = payload["indicative_range_usd"]
        textbox(slide, 0.7, 2.1, 11.5, 0.5,
                f"Indicative range: USD {low:,.0f} – {high:,.0f}", size=18)
        textbox(slide, 0.7, 2.8, 11.5, 1.0,
                "The indicative range is not a quantified exposure: its parameters "
                "are uncalibrated judgement inputs, so no point estimate is given.",
                size=11, color=TEXT_BODY)

    # Findings slides, six per slide
    findings = payload["findings"]
    for start in range(0, len(findings), 6):
        slide = presentation.slides.add_slide(blank)
        textbox(slide, 0.7, 0.4, 11.5, 0.6,
                f"Findings ({start + 1}–{min(start + 6, len(findings))} of "
                f"{len(findings)})", size=26, bold=True)
        top = 1.3
        for finding in findings[start:start + 6]:
            label = finding.get("finding") or finding.get("title", "Finding")
            textbox(slide, 0.7, top, 11.9, 0.85,
                    f"[{finding.get('severity', 'Unrated')}] {label}", size=12)
            top += 0.9

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def export(session, report_id, payload, fmt="docx"):
    """Render and release. Blocked unless sign-off covers these exact bytes."""
    principal = require_principal()
    renderers = {"docx": render_docx, "pptx": render_pptx}
    if fmt not in renderers:
        raise ExportError(f"Unsupported format {fmt!r}. Use one of: docx, pptx.")

    content = renderers[fmt](payload)
    signoff.require_releasable(session, report_id, content)

    audit.record(
        session, principal.username, "report.exported",
        entity_type="report_signoff", entity_id=report_id,
        deal_id=payload["deal_id"],
        detail={"format": fmt, "bytes": len(content),
                "content_sha256": signoff.content_hash(content),
                "findings": len(payload["findings"])},
    )
    return content


def prepare_for_signoff(session, report_id, payload, fmt="docx"):
    """Render the candidate deliverable and open sign-off against its bytes."""
    renderers = {"docx": render_docx, "pptx": render_pptx}
    content = renderers[fmt](payload)
    digest = signoff.open_signoff(session, report_id, payload["deal_id"], content)
    return content, digest
