from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_BG = RGBColor(30, 30, 35)
DARK_GRAY = RGBColor(55, 55, 60)
MED_GRAY = RGBColor(130, 130, 130)
LIGHT_GRAY = RGBColor(220, 220, 220)
ORANGE = RGBColor(255, 90, 0)
ORANGE_DARK = RGBColor(210, 70, 0)
RED_ACCENT = RGBColor(200, 40, 40)
NEAR_WHITE = RGBColor(248, 248, 248)
BOX_BG = RGBColor(252, 252, 252)
BORDER_GRAY = RGBColor(210, 210, 210)
TEXT_DARK = RGBColor(40, 40, 45)
TEXT_BODY = RGBColor(60, 60, 65)

FONT = "Calibri"
FONT_LIGHT = "Calibri Light"


def rect(slide, l, t, w, h, fill, line=None, line_w=Pt(0.75)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, fill, line=None, line_w=Pt(0.75)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb


def add_bullets(slide, l, t, w, h, items, sz=10, color=TEXT_BODY, prefix="☐  ", spacing=Pt(4), bold_items=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    for i, item in enumerate(items):
        p = tf.add_paragraph()
        p.text = f"{prefix}{item}"
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_before = spacing
        p.space_after = Pt(1)
        if bold_items and i in bold_items:
            p.font.bold = True
    return tb


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1: TITLE (matches reference pattern exactly)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Dark full background
rect(slide1, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BG)

# Subtle gradient overlay on right half (simulating image area)
rect(slide1, Inches(6.5), Inches(0), Inches(6.833), Inches(7.5), RGBColor(40, 40, 48))
rect(slide1, Inches(8), Inches(0), Inches(5.333), Inches(7.5), RGBColor(50, 50, 58))

# Orange swoosh ribbons (diagonal parallelogram-like shapes)
# Upper swoosh
s1 = slide1.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(4.5), Inches(3.0), Inches(5.5), Inches(0.7))
s1.fill.solid()
s1.fill.fore_color.rgb = ORANGE
s1.line.fill.background()
s1.rotation = -2.0

# Lower swoosh
s2 = slide1.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(7), Inches(3.8), Inches(4.5), Inches(0.9))
s2.fill.solid()
s2.fill.fore_color.rgb = ORANGE
s2.line.fill.background()
s2.rotation = -3.0

# PwC logo text
txt(slide1, Inches(0.6), Inches(0.4), Inches(2), Inches(0.7),
    "pwc", sz=28, bold=True, color=WHITE, font=FONT_LIGHT)

# Main title
txt(slide1, Inches(0.6), Inches(1.8), Inches(9), Inches(1.2),
    "PwC India Data & Analytics", sz=44, bold=True, color=WHITE, font=FONT_LIGHT)

# Subtitle
txt(slide1, Inches(0.6), Inches(3.2), Inches(9), Inches(0.9),
    "Gen AI & Agentic AI Use Cases", sz=34, bold=True, color=WHITE, font=FONT_LIGHT)

# Date
txt(slide1, Inches(0.6), Inches(6.7), Inches(3), Inches(0.4),
    "August 2026", sz=13, color=MED_GRAY, font=FONT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2: CONTENT (matches reference 4-quadrant layout exactly)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
rect(slide2, Inches(0), Inches(0), Inches(13.333), Inches(7.5), WHITE)

# ── LEFT PANEL: Dark strip with image area ──
LEFT_W = Inches(3.1)
rect(slide2, Inches(0), Inches(0), LEFT_W, Inches(7.5), DARK_BG)

# Category tag (teal/dark blue banner)
tag = rrect(slide2, Inches(0.15), Inches(0.3), Inches(2.8), Inches(0.45), RGBColor(20, 80, 100))
txt(slide2, Inches(0.15), Inches(0.32), Inches(2.8), Inches(0.42),
    "ESG & SUSTAINABILITY", sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Simulated image area (dark textured rectangle)
rect(slide2, Inches(0), Inches(0.9), LEFT_W, Inches(2.2), RGBColor(25, 25, 30))
# Add some faint visual elements to suggest analytics/data imagery
for i in range(5):
    bar = rect(slide2, Inches(0.3 + i * 0.5), Inches(2.3 - i * 0.15), Inches(0.15), Inches(0.4 + i * 0.12),
               RGBColor(60 + i * 8, 60 + i * 5, 70 + i * 6))
    bar.fill.fore_color.rgb = RGBColor(60 + i * 10, 70 + i * 8, 90 + i * 6)

# ── KEY HIGHLIGHTS BOX ──
KH_TOP = Inches(3.2)
KH_H = Inches(4.15)
rect(slide2, Inches(0.1), KH_TOP, Inches(2.9), KH_H, RGBColor(245, 245, 245), RGBColor(180, 180, 180))

txt(slide2, Inches(0.2), KH_TOP + Inches(0.1), Inches(2.7), Inches(0.4),
    "Key Highlights", sz=16, bold=True, color=BLACK)

# Key highlight bullet items with bold lead-ins
kh_tb = slide2.shapes.add_textbox(Inches(0.2), KH_TOP + Inches(0.5), Inches(2.7), Inches(3.6))
kh_tf = kh_tb.text_frame
kh_tf.word_wrap = True
kh_tf.paragraphs[0].text = ""

kh_items = [
    ("6-Agent Orchestration Pipeline", " powers the entire workflow end-to-end."),
    ("9 Connector Types", " for data ingestion: File Upload, Google Sheets, REST API, SQL DB, AWS S3, BigQuery, GCS, Azure Blob, Delta Lake."),
    ("23 Canonical Tables", " (auto-detected): Companies, ESG Metrics, Compliance, Risks, Financials, Supply Chain, Benchmarks."),
    ("9 Compliance Frameworks", " covered simultaneously: BRSR, CSRD, GRI, TCFD, SASB, GDPR, SEC, DPDP, SOX."),
    ("Deal-Ready Governance:", " Role-based access (Admin, Manager, Analyst, Viewer) + audit trail + quality scoring + evidence traceability."),
]

for bold_text, normal_text in kh_items:
    p = kh_tf.add_paragraph()
    p.space_before = Pt(5)
    p.space_after = Pt(1)
    run_b = p.add_run()
    run_b.text = f"•  {bold_text}"
    run_b.font.size = Pt(9)
    run_b.font.bold = True
    run_b.font.color.rgb = TEXT_DARK
    run_b.font.name = FONT
    run_n = p.add_run()
    run_n.text = normal_text
    run_n.font.size = Pt(9)
    run_n.font.bold = False
    run_n.font.color.rgb = TEXT_BODY
    run_n.font.name = FONT


# ── RIGHT AREA: Title + 4 Quadrants ──
RIGHT_X = Inches(3.3)
RIGHT_W = Inches(9.9)

# Main title at top
txt(slide2, RIGHT_X + Inches(0.2), Inches(0.2), Inches(9.5), Inches(0.6),
    "AI-Driven ESG Due Diligence & Data Intelligence Platform",
    sz=22, bold=True, color=BLACK, font=FONT)

# Quadrant dimensions
Q_MARGIN = Inches(0.15)
Q_W = Inches(4.75)
Q_H = Inches(3.0)
Q_ROW1_Y = Inches(0.9)
Q_ROW2_Y = Inches(4.05)
Q_COL1_X = RIGHT_X + Inches(0.1)
Q_COL2_X = RIGHT_X + Inches(5.05)

# ── QUADRANT 1: BACKGROUND (top-left) ──
rect(slide2, Q_COL1_X, Q_ROW1_Y, Q_W, Q_H, BOX_BG, BORDER_GRAY)
# Header bar
rect(slide2, Q_COL1_X, Q_ROW1_Y, Q_W, Inches(0.4), NEAR_WHITE, BORDER_GRAY)
txt(slide2, Q_COL1_X + Inches(0.15), Q_ROW1_Y + Inches(0.05), Inches(3), Inches(0.35),
    "Background", sz=13, bold=True, color=TEXT_DARK)
# Icon placeholder
icon_bg = rrect(slide2, Q_COL1_X + Q_W - Inches(0.55), Q_ROW1_Y + Inches(0.05), Inches(0.35), Inches(0.3), NEAR_WHITE, BORDER_GRAY)
txt(slide2, Q_COL1_X + Q_W - Inches(0.55), Q_ROW1_Y + Inches(0.04), Inches(0.35), Inches(0.3),
    "☰", sz=12, color=MED_GRAY, align=PP_ALIGN.CENTER)

bg_items = [
    "The ESG Data Platform is an agentic AI platform that transforms ESG due diligence from a periodic manual task into an always-on intelligent workflow.",
    "It targets M&A deal teams performing acquisition due diligence, IPO readiness, and growth equity analysis across India, EU, and US jurisdictions.",
    "Built on Streamlit + Claude AI (Anthropic) on Google Vertex AI, it replaces fragmented spreadsheets with a unified 6-agent pipeline.",
]
add_bullets(slide2, Q_COL1_X + Inches(0.1), Q_ROW1_Y + Inches(0.5), Q_W - Inches(0.2), Q_H - Inches(0.55),
            bg_items, sz=9, spacing=Pt(6))


# ── QUADRANT 2: CHALLENGES (top-right) ──
rect(slide2, Q_COL2_X, Q_ROW1_Y, Q_W, Q_H, BOX_BG, BORDER_GRAY)
rect(slide2, Q_COL2_X, Q_ROW1_Y, Q_W, Inches(0.4), NEAR_WHITE, BORDER_GRAY)
txt(slide2, Q_COL2_X + Inches(0.15), Q_ROW1_Y + Inches(0.05), Inches(3), Inches(0.35),
    "Challenges", sz=13, bold=True, color=TEXT_DARK)
icon_ch = rrect(slide2, Q_COL2_X + Q_W - Inches(0.55), Q_ROW1_Y + Inches(0.05), Inches(0.35), Inches(0.3), NEAR_WHITE, BORDER_GRAY)
txt(slide2, Q_COL2_X + Q_W - Inches(0.55), Q_ROW1_Y + Inches(0.04), Inches(0.35), Inches(0.3),
    "⚠", sz=12, color=ORANGE, align=PP_ALIGN.CENTER)

ch_items = [
    "Data silos across sustainability reports, PDFs, regulatory filings, supplier assessments. Manual reconciliation is slow & error-prone.",
    "Framework Complexity: BRSR + CSRD + GRI + TCFD + SASB + GDPR + SEC + DPDP + SOX — different metrics, thresholds, cadences.",
    "Reporting Latency: Traditional due diligence takes weeks per deal. Backward-looking, not actionable in real time.",
    "Greenwashing Risk: Self-reported metrics diverge from operational data. No systematic integrity gap detection.",
    "No Quantified Risk Exposure: Deal teams need hard numbers — financial exposure by ESG pillar, not subjective scores.",
]
add_bullets(slide2, Q_COL2_X + Inches(0.1), Q_ROW1_Y + Inches(0.5), Q_W - Inches(0.2), Q_H - Inches(0.55),
            ch_items, sz=9, spacing=Pt(4))


# ── QUADRANT 3: APPROACH (bottom-left) ──
rect(slide2, Q_COL1_X, Q_ROW2_Y, Q_W, Q_H, BOX_BG, BORDER_GRAY)
rect(slide2, Q_COL1_X, Q_ROW2_Y, Q_W, Inches(0.4), NEAR_WHITE, BORDER_GRAY)
txt(slide2, Q_COL1_X + Inches(0.15), Q_ROW2_Y + Inches(0.05), Inches(3), Inches(0.35),
    "Approach", sz=13, bold=True, color=TEXT_DARK)
icon_ap = rrect(slide2, Q_COL1_X + Q_W - Inches(0.55), Q_ROW2_Y + Inches(0.05), Inches(0.35), Inches(0.3), NEAR_WHITE, BORDER_GRAY)
txt(slide2, Q_COL1_X + Q_W - Inches(0.55), Q_ROW2_Y + Inches(0.04), Inches(0.35), Inches(0.3),
    "⚙", sz=12, color=MED_GRAY, align=PP_ALIGN.CENTER)

ap_items = [
    "Registration Agent: 9 connectors · auto-schema detection · AI column mapping to 23 canonical tables.",
    "Metric Analysis Agent: YoY trends · CAGR · intensity metrics · anomaly detection · target progress tracking.",
    "Regulatory Tracker Agent: Monitors BRSR · CSRD · GRI · TCFD · SASB · GDPR · SEC · DPDP · SOX for compliance gaps.",
    "Benchmarking Agent: Tiered peer comparison · percentile ranking · performance classification across 19+ companies.",
    "Risk & Opportunity Agent: 5×5 risk matrix · financial exposure by pillar · deal recommendations · opportunity cards.",
    "Review & Report Agent: Quality scoring · conflict detection · governance validation · executive report generation.",
]
add_bullets(slide2, Q_COL1_X + Inches(0.1), Q_ROW2_Y + Inches(0.5), Q_W - Inches(0.2), Q_H - Inches(0.55),
            ap_items, sz=9, spacing=Pt(3))


# ── QUADRANT 4: IMPACT (bottom-right) ──
rect(slide2, Q_COL2_X, Q_ROW2_Y, Q_W, Q_H, BOX_BG, BORDER_GRAY)
# Red header bar (matching reference pattern)
rect(slide2, Q_COL2_X, Q_ROW2_Y, Q_W, Inches(0.4), RED_ACCENT)
txt(slide2, Q_COL2_X + Inches(0.15), Q_ROW2_Y + Inches(0.05), Inches(3), Inches(0.35),
    "Impact", sz=13, bold=True, color=WHITE)

# Impact transformation metrics
imp_transforms = [
    "Due Diligence Cycle: Weeks/Months → Hours — Real-time pipeline",
    "Framework Coverage: Single framework → 9 frameworks simultaneously",
    "Risk Assessment: Subjective → Quantified financial exposure by E/S/G pillar",
    "Data Quality: Manual checks → AI-powered quality scoring + anomaly detection",
    "Governance: Ad-hoc review → Role-based access + audit trail + conflict detection",
]
add_bullets(slide2, Q_COL2_X + Inches(0.1), Q_ROW2_Y + Inches(0.45), Q_W - Inches(0.2), Inches(1.3),
            imp_transforms, sz=9, spacing=Pt(3))

# Key benefits sub-section
txt(slide2, Q_COL2_X + Inches(0.15), Q_ROW2_Y + Inches(1.6), Inches(2), Inches(0.25),
    "Key Benefits:", sz=10, bold=True, color=TEXT_DARK)

benefits = [
    "Cycles: 80% faster.",
    "Agents: 6 AI agents orchestrated.",
    "Compliance Gaps: Auto-detected with remediation plans.",
    "Frameworks: 9 simultaneous.",
    "Peer Benchmarks: 19+ companies · 5 countries.",
    "Deals Supported: M&A, Due Diligence, IPO, Growth Equity.",
]
add_bullets(slide2, Q_COL2_X + Inches(0.1), Q_ROW2_Y + Inches(1.8), Q_W - Inches(0.2), Inches(1.0),
            benefits, sz=9, spacing=Pt(2))

# Next Steps sub-section
txt(slide2, Q_COL2_X + Inches(0.15), Q_ROW2_Y + Inches(2.55), Inches(2), Inches(0.2),
    "Next Steps:", sz=10, bold=True, color=TEXT_DARK)

next_steps = [
    "Expand connector ecosystem for real-time enterprise data ingestion (SAP, Workday).",
    "From compliance checkbox to strategic intelligence — Predictive risk scoring · Quantified deal impact · Actionable recommendations.",
]
add_bullets(slide2, Q_COL2_X + Inches(0.1), Q_ROW2_Y + Inches(2.7), Q_W - Inches(0.2), Inches(0.6),
            next_steps, sz=9, spacing=Pt(2))

# Page number
txt(slide2, Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.3),
    "2", sz=10, color=MED_GRAY, align=PP_ALIGN.RIGHT)


# ━━━ SAVE ━━━
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "PwC_ESG_Platform_UseCase.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
