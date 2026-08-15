"""Report export gates: reviewed exposures, publishable benchmarks, sign-off."""

import pytest

from esg.benchmarks import provenance
from esg.db.models import CompanyFinancials, PeerBenchmarkData
from esg.db.scope import Principal, bind_principal
from esg.deal import red_flag_report, signoff
from esg.methodology import exposure

FINDINGS = [
    {"severity": "High", "finding": "Group Scope 1 does not reconcile to facilities",
     "why_it_matters": "Suggests an excluded site.",
     "reported_citation": "Sustainability Report 2024.pdf, p.42"},
    {"severity": "Critical", "finding": "Unpaid environmental penalty outstanding",
     "why_it_matters": "Crystallised liability.",
     "citation": "Consent order.pdf, p.3"},
    {"severity": "Medium", "finding": "Attrition above sector norm"},
]


@pytest.fixture
def deal(session, admin, deal_setup):
    with bind_principal(admin):
        session.add(CompanyFinancials(company_id="C1", reporting_year=2024,
                                      annual_revenue=200_000_000,
                                      reporting_currency="USD"))
        session.commit()
    return {
        **deal_setup,
        "manager": Principal("u-m", "deal_manager", "Manager", {"D1": "Owner"}),
    }


def test_payload_orders_findings_by_severity(session, deal):
    with bind_principal(deal["analyst"]):
        payload = red_flag_report.build_payload(
            session, "D1", "C1", FINDINGS, reporting_year=2024
        )
    assert [f["severity"] for f in payload["findings"]] == ["Critical", "High", "Medium"]
    assert payload["methodology_version"] == exposure.METHODOLOGY_VERSION


def test_unreviewed_exposure_blocks_the_report(session, deal):
    with bind_principal(deal["analyst"]):
        run, _ = exposure.quantify(
            session, {"finding_id": "F1", "severity": "High", "impact_score": 4},
            company_id="C1", deal_id="D1",
        )
        session.commit()
        with pytest.raises(exposure.ReviewRequired):
            red_flag_report.build_payload(session, "D1", "C1", FINDINGS,
                                          exposures=[run])


def test_illustrative_benchmarks_cannot_leave_the_building(session, deal, admin):
    with bind_principal(admin):
        session.add_all([
            PeerBenchmarkData(
                benchmark_record_id=f"BM{i}", peer_company_name=f"Fictional {i}",
                metric_code="ENV_SCOPE1", metric_value=100.0,
                provenance=provenance.ILLUSTRATIVE,
            ) for i in range(10)
        ])
        session.commit()

    cohort = provenance.cohort_provenance(session, "ENV_SCOPE1")
    with bind_principal(deal["analyst"]):
        with pytest.raises(provenance.BenchmarkProvenanceError, match="red-flag report"):
            red_flag_report.build_payload(session, "D1", "C1", FINDINGS,
                                          benchmark_cohorts=[cohort])


def test_indicative_exposure_is_labelled_not_summed_into_quantified(session, deal):
    with bind_principal(deal["analyst"]):
        run, _ = exposure.quantify(
            session, {"finding_id": "F2", "severity": "High", "impact_score": 4},
            company_id="C1", deal_id="D1",
        )
        session.commit()
    with bind_principal(deal["manager"]):
        exposure.review(session, run.exposure_run_id, "checked")
        session.commit()
        payload = red_flag_report.build_payload(session, "D1", "C1", FINDINGS,
                                                exposures=[run])

    assert payload["quantified_total_usd"] == 0.0
    assert payload["indicative_range_usd"] is not None
    assert payload["exposures"][0]["presentation"]["disclosure"] == exposure.INDICATIVE


def test_docx_renders_and_carries_citations(session, deal):
    with bind_principal(deal["analyst"]):
        payload = red_flag_report.build_payload(session, "D1", "C1", FINDINGS)
        content = red_flag_report.render_docx(payload)
    assert content[:2] == b"PK"  # a valid OOXML package

    import io

    import docx

    document = docx.Document(io.BytesIO(content))
    text = "\n".join(p.text for p in document.paragraphs)
    assert "Consent order.pdf, p.3" in text
    assert "Red Flag Report" in text


def test_pptx_renders_with_run_level_formatting(session, deal):
    """The original deck stored formatting in paragraph defaults, which other
    renderers ignore. Exported decks must put it on the run."""
    import io

    from pptx import Presentation

    with bind_principal(deal["analyst"]):
        payload = red_flag_report.build_payload(session, "D1", "C1", FINDINGS)
        content = red_flag_report.render_pptx(payload)

    presentation = Presentation(io.BytesIO(content))
    sized_runs = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size is not None:
                        sized_runs += 1
    assert sized_runs > 5, "font sizes must live on runs, not paragraph defaults"


def test_export_is_blocked_until_signed(session, deal):
    with bind_principal(deal["analyst"]):
        payload = red_flag_report.build_payload(session, "D1", "C1", FINDINGS)
        with pytest.raises(signoff.NotReleasable):
            red_flag_report.export(session, "RPT-X", payload)


def test_full_signoff_then_export_succeeds(session, deal):
    with bind_principal(deal["analyst"]):
        payload = red_flag_report.build_payload(session, "D1", "C1", FINDINGS)
        content, digest = red_flag_report.prepare_for_signoff(session, "RPT-Y", payload)
        session.commit()

    with bind_principal(deal["manager"]):
        signoff.sign(session, "RPT-Y", content, "Approved", "Reviewed to evidence")
        session.commit()

    with bind_principal(deal["analyst"]):
        exported = red_flag_report.export(session, "RPT-Y", payload)
    assert signoff.content_hash(exported) == digest


def test_editing_the_report_after_signoff_blocks_export(session, deal):
    with bind_principal(deal["analyst"]):
        payload = red_flag_report.build_payload(session, "D1", "C1", FINDINGS)
        content, _ = red_flag_report.prepare_for_signoff(session, "RPT-Z", payload)
        session.commit()
    with bind_principal(deal["manager"]):
        signoff.sign(session, "RPT-Z", content, "Approved")
        session.commit()

    with bind_principal(deal["analyst"]):
        payload["findings"].append({"severity": "Critical", "finding": "Added later"})
        with pytest.raises(signoff.NotReleasable, match="does not match what was signed"):
            red_flag_report.export(session, "RPT-Z", payload)


def test_outstanding_irs_appear_in_the_report(session, deal):
    from esg.deal import information_requests as irs

    with bind_principal(deal["analyst"]):
        irs.raise_request(session, "D1", "C1", "Provide Scope 1 workings",
                          priority="High")
        session.commit()
        payload = red_flag_report.build_payload(session, "D1", "C1", FINDINGS)
        content = red_flag_report.render_docx(payload)

    import io

    import docx

    text = "\n".join(p.text for p in docx.Document(io.BytesIO(content)).paragraphs)
    assert "Provide Scope 1 workings" in text
    assert payload["outstanding_information_requests"]["count"] == 1
